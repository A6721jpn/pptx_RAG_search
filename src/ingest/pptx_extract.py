import pptx
from typing import List, Dict
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

class PptxExtractor:
    """Extracts text and metadata from PowerPoint files."""
    
    def extract(self, file_path: Path) -> List[Dict]:
        """
        Extract text and notes from a PPTX file.
        
        Args:
            file_path: Path to the .pptx file
            
        Returns:
            List of dictionaries containing extracted data for each slide.
            Each dict has: slide_no, text_raw, notes_raw
        """
        try:
            prs = pptx.Presentation(file_path)
            slides_data = []

            for i, slide in enumerate(prs.slides):
                slide_no = i + 1
                text_parts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    # Check for text frame
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                text_parts.append(text)
                    # Use generic .text if available (e.g. some shapes)
                    elif hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)
                    
                    # TODO: Handle tables if necessary (python-pptx supports it)

                text_raw = "\n".join(text_parts)
                
                # Extract notes
                notes_raw = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_raw = notes_slide.notes_text_frame.text.strip()

                slides_data.append({
                    "slide_no": slide_no,
                    "text_raw": text_raw,
                    "notes_raw": notes_raw
                })
                
            return slides_data
            
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            raise
